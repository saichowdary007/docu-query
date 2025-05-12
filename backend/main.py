import os
import io
import tempfile
import uuid
import logging
from typing import List, Dict, Any, Optional
import pandas as pd
import pdfplumber # type: ignore
import docx # type: ignore
import pptx # type: ignore
import re # Added for sanitization

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware

import google.generativeai as genai # type: ignore
from langchain_core.documents import Document
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.prompts import ChatPromptTemplate
from dotenv import load_dotenv
from langchain.text_splitter import RecursiveCharacterTextSplitter, MarkdownTextSplitter

# RAG core and response generator
from rag_core import embed, retrieve, register_duckdb_table, get_shared_db_connection, execute_sql_query_rag
from response_generator import process_query_for_direct_response

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Configure Google Generative AI
api_key = os.getenv("GOOGLE_API_KEY")
if api_key:
    try:
        genai.configure(api_key=api_key)
        logger.info(f"Google API Key configured: {api_key[:5]}...")
    except Exception as e:
        logger.error(f"Error configuring Google API Key for genai package: {e}")
else:
    logger.warning("Google API Key not configured. LLM features will be limited.")

# Initialize FastAPI app
app = FastAPI(title="DocuQuery API")

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # Allows all origins
    allow_credentials=True,
    allow_methods=["*"], # Allows all methods
    allow_headers=["*"], # Allows all headers
)

# Load system prompt
try:
    DETAILED_SYSTEM_PROMPT = """
You are DocuQuery, an advanced document intelligence system. Your primary goal is to provide precise, direct, and comprehensive answers based *solely* on the information contained within the documents provided to you in the context.

CORE DIRECTIVES:
1.  **Accuracy is Paramount:** Only state facts and information explicitly present in the provided document excerpts. Do NOT infer, speculate, or add external knowledge. If the information isn't in the context, clearly state that.
2.  **Direct Answers First:** Begin your response with a direct answer to the user's question. Follow up with relevant supporting details, explanations, and elaborations from the documents.
3.  **Comprehensive Synthesis:** If multiple document excerpts are relevant, synthesize them into a coherent and complete answer. Avoid simply listing snippets.
4.  **Clarity and Conciseness:** Use clear, straightforward language. Format your responses for readability (e.g., bullet points for lists, bolding for emphasis on key terms or results).
5.  **Handle Ambiguity:** If the user's query is ambiguous or if the provided context is insufficient to fully answer, ask clarifying questions or state what information is missing.

DOCUMENT EXPERTISE & RESPONSE STRATEGY:

* **Tabular Data (CSV, Excel, Sheets):**
    * When asked about data in tables, focus on extracting specific values, performing calculations if requested (e.g., sums, averages based on provided data), and identifying trends or patterns *visible within the given rows/schema*.
    * If the user's query implies a complex data operation (e.g., "Show me all sales from Q1 filtered by region X and sorted by date"), and you have been provided with an SQL query and its results, explain the results clearly. If only schema or row snippets are available, you can suggest that a more complex SQL query on the full dataset might be needed for a complete answer, but do not attempt to generate complex SQL yourself unless explicitly instructed by the system in a special way (e.g., via a `direct_response_payload`).
    * Refer to table names or column names if they are known from the context and relevant to the answer.

* **Text Documents (PDF, DOCX, TXT, MD):**
    * Extract specific facts, definitions, explanations, or summaries as requested.
    * Identify key entities (like people, organizations, locations, dates) if relevant to the query and present in the text.
    * Pay attention to document structure if evident from the context (e.g., headings, sections mentioned in metadata).

* **Presentations (PPTX):**
    * Extract information from slide content. Refer to slide numbers or key phrases from the slide if it helps clarify the source of the information within the presentation context.

BEST PRACTICES FOR HIGH-QUALITY ANSWERS:
* **Thorough Contextual Analysis:** Carefully read and understand all provided document excerpts before formulating your answer.
* **No External Information:** Your knowledge is strictly limited to the documents given in the current context. Do not use any prior knowledge or information from outside these documents.
* **Acknowledge Limitations:** If the answer cannot be found or fully determined from the provided context, explicitly state this. For example, "Based on the provided document excerpts, I cannot determine X..." or "The documents do not contain information about Y."
* **Focus on User's Goal:** Understand the intent behind the user's question and provide the most helpful and relevant information from the documents.
* **Numerical Precision:** When extracting numerical data, ensure accuracy and include units if specified in the text.
* **Comparisons:** If asked to compare items, create a structured comparison based *only* on attributes found in the provided documents for those items.
* **Extraction, Not Interpretation:** Stick to extracting and presenting information. Avoid subjective interpretations or opinions.
"""
    SYSTEM_PROMPT = DETAILED_SYSTEM_PROMPT
    logger.info("Detailed system prompt loaded successfully")
except Exception as e:
    logger.error(f"Error loading system prompt: {e}. Using a basic default.")
    SYSTEM_PROMPT = "You are DocuQuery, an expert document analysis assistant. Provide answers based only on the document context."


# Initialize LLM
llm: Any 
if api_key:
    try:
        llm = ChatGoogleGenerativeAI(
            model="gemini-pro", 
            temperature=0.1, # Slightly lower temperature for more factual responses
            convert_system_message_to_human=True,
            google_api_key=api_key
        )
        logger.info("Google Gemini LLM initialized successfully (gemini-pro).")
    except Exception as e:
        logger.error(f"Error initializing Google Gemini LLM: {e}. Using fallback FakeListLLM.")
        from langchain_core.language_models.fake import FakeListLLM
        llm = FakeListLLM(responses=["I'm sorry, but I cannot process your request. The primary language model is unavailable."])
else:
    from langchain_core.language_models.fake import FakeListLLM
    llm = FakeListLLM(responses=["LLM functionality is disabled as no API key is provided."])
    logger.info("Using FakeListLLM as LLM is not configured (no API key).")


# File upload directory
UPLOAD_DIR = "./uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

# --- Enhanced Document Loaders with Improved Chunking ---

def load_csv_xlsx(file_path: str, file_id: str, original_filename: str) -> List[Document]:
    """Load CSV or XLSX file into documents and register as DuckDB table"""
    documents = []
    df = None
    logger.info(f"Attempting to load tabular file: {original_filename} from path: {file_path}")
    try:
        if file_path.endswith('.csv'):
            try:
                df = pd.read_csv(file_path)
                logger.info(f"Successfully read CSV '{original_filename}' with pandas.")
            except pd.errors.ParserError as pe:
                logger.warning(f"Pandas ParserError for CSV '{original_filename}': {pe}. Trying with different encoding/sep.")
                try:
                    df = pd.read_csv(file_path, encoding='latin1')
                    logger.info(f"Successfully read CSV '{original_filename}' with latin1 encoding.")
                except Exception as e_latin1:
                    logger.warning(f"Failed to read CSV '{original_filename}' with latin1: {e_latin1}. Trying with sep=';'.")
                    try:
                        df = pd.read_csv(file_path, sep=';')
                        logger.info(f"Successfully read CSV '{original_filename}' with sep=';'.")
                    except Exception as e_sep:
                        logger.error(f"All CSV reading attempts failed for '{original_filename}': {e_sep}", exc_info=True)
                        raise  # Re-raise to be caught by the outer try-except
            except Exception as e_csv:
                logger.error(f"Generic error loading CSV '{original_filename}': {e_csv}", exc_info=True)
                raise # Re-raise
        else: # Excel files
            engine = None
            if file_path.endswith('.xlsx'): engine = 'openpyxl'
            elif file_path.endswith('.xls'): engine = 'xlrd'
            elif file_path.endswith('.xlsb'): engine = 'pyxlsb'
            
            logger.info(f"Attempting to read Excel file '{original_filename}' with engine: {engine if engine else 'default pandas'}")
            try:
                df = pd.read_excel(file_path, engine=engine)
                logger.info(f"Successfully read Excel file '{original_filename}' with engine: {engine if engine else 'default pandas'}.")
            except Exception as xls_error:
                logger.warning(f"Error with primary Excel engine for '{original_filename}' (engine: {engine}): {xls_error}", exc_info=True)
                if engine == 'xlrd' and "CompDoc" in str(xls_error): # Compound document error often means it's not a true BIFF .xls
                    logger.warning(f"'{original_filename}' might not be a standard BIFF .xls file. xlrd error: {xls_error}")
                
                logger.info(f"Trying fallback Excel reading for '{original_filename}' (pandas default engine).")
                try:
                    df = pd.read_excel(file_path) # Fallback to default
                    logger.info(f"Successfully read Excel file '{original_filename}' with pandas default engine.")
                except Exception as final_xls_error:
                    logger.error(f"All Excel reading methods failed for '{original_filename}': {final_xls_error}", exc_info=True)
                    raise # Re-raise
        
        if df is None or df.empty:
            logger.warning(f"File {original_filename} is empty or failed to load into DataFrame after all attempts.")
            return [] # Return empty list if df is still None or empty
        
        logger.info(f"Loaded tabular file '{original_filename}' with {len(df)} rows, {len(df.columns)} columns.")

        table_name_prefix = os.path.splitext(original_filename)[0]
        safe_prefix = re.sub(r'[^a-zA-Z0-9_]', '_', table_name_prefix).lower()
        safe_file_id_suffix = file_id.replace('-', '')[:8] 
        duckdb_table_name_base = f"{safe_prefix}_{safe_file_id_suffix}"
        
        table_name = register_duckdb_table(df, duckdb_table_name_base)
        if not table_name:
            logger.error(f"DuckDB table registration failed for '{original_filename}'")
            return [] 

        # Schema Document
        schema_info_parts = [f"Table Name: \"{table_name}\"", "Columns (with types):"]
        try:
            db_conn_temp = get_shared_db_connection()
            cols_info_df = db_conn_temp.execute(f"PRAGMA table_info(\"{table_name}\")").fetchdf()
            for _, row in cols_info_df.iterrows():
                schema_info_parts.append(f"- \"{row['name']}\": {row['type']}")
            schema_info_parts.append(f"Total Rows: {len(df)}")
            schema_info = "\n".join(schema_info_parts)
        except Exception as e:
            logger.warning(f"Could not fetch detailed schema for {table_name}: {e}")
            schema_info = f"Table: \"{table_name}\"\nColumns: {', '.join(df.columns.tolist())}\nRows: {len(df)}"

        documents.append(Document(
            page_content=schema_info,
            metadata={
                "source": original_filename, "file_id": file_id, "file_type": "tabular_schema",
                "table_name": table_name, "is_schema": True, "doc_id_prefix": f"{file_id}_schema"
            }
        ))

        # Statistics Document
        try:
            numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
            if numeric_cols:
                stats = df[numeric_cols].describe().to_string()
                documents.append(Document(
                    page_content=f"Statistical summary for table \"{table_name}\":\n{stats}",
                    metadata={
                        "source": original_filename, "file_id": file_id, "file_type": "tabular_stats",
                        "table_name": table_name, "is_stats": True, "doc_id_prefix": f"{file_id}_stats"
                    }
                ))
        except Exception as stats_error:
            logger.warning(f"Could not generate statistics for '{original_filename}': {stats_error}")

        # Row-level documents
        for i, row_data in df.iterrows():
            row_dict_clean = {str(k): (str(v) if pd.notna(v) else None) for k, v in row_data.to_dict().items()}
            content = f"Row {i+1} from table \"{table_name}\" ({original_filename}): " + \
                      "; ".join([f"\"{col}\": \"{val}\"" for col, val in row_dict_clean.items() if val is not None])
            documents.append(Document(
                page_content=content,
                metadata={
                    "row_index": i + 1, "source": original_filename, "file_id": file_id,
                    "table_name": table_name, "file_type": "tabular_row", "is_tabular": True,
                    "doc_id_prefix": f"{file_id}_row_{i+1}",
                    **row_dict_clean 
                }
            ))
        
        logger.info(f"Created {len(documents)} documents from tabular file '{original_filename}' for table '{table_name}'")
        return documents
    except Exception as e: # Catch-all for the entire function
        logger.error(f"Overall error in load_csv_xlsx for {original_filename}: {e}", exc_info=True)
        return []


def load_pdf(file_path: str, file_id: str, original_filename: str) -> List[Document]:
    """Load PDF file into documents with improved chunking and OCR attempt."""
    documents = []
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200, add_start_index=True)
    try:
        with pdfplumber.open(file_path) as pdf:
            logger.info(f"PDF '{original_filename}' has {len(pdf.pages)} pages")
            full_text_for_ocr_check = ""
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                page_number = i + 1
                ocr_attempted = False
                if not text or len(text.strip()) < 20: # If little text, try OCR
                    logger.info(f"Page {page_number} of {original_filename} has minimal direct text. Attempting OCR.")
                    ocr_attempted = True
                    try:
                        import pytesseract # type: ignore
                        # from PIL import Image # Not strictly needed if page.to_image works
                        pil_image = page.to_image(resolution=200).original 
                        text_from_ocr = pytesseract.image_to_string(pil_image)
                        if text_from_ocr and text_from_ocr.strip():
                             logger.info(f"Extracted text using OCR for page {page_number} of {original_filename}")
                             text = text_from_ocr # Use OCR text
                        else:
                            logger.info(f"OCR for page {page_number} of {original_filename} yielded no significant text.")
                            if not text: text = "" # Ensure text is empty string if OCR also fails
                    except ImportError:
                        logger.warning("pytesseract not installed. OCR for PDF images skipped.")
                        if not text: text = "" 
                    except Exception as ocr_error:
                        logger.warning(f"OCR failed for page {page_number} of {original_filename}: {ocr_error}")
                        if not text: text = "" 
                
                if text and text.strip():
                    full_text_for_ocr_check += text + "\n"
                    page_chunks = text_splitter.split_text(text)
                    for chunk_index, chunk in enumerate(page_chunks):
                        doc = Document(
                            page_content=chunk,
                            metadata={
                                "page": page_number,
                                "chunk_on_page": chunk_index + 1,
                                "source": original_filename,
                                "file_id": file_id,
                                "file_type": "pdf_chunk",
                                "ocr_attempted": ocr_attempted,
                                "doc_id_prefix": f"{file_id}_page_{page_number}_chunk_{chunk_index+1}"
                            }
                        )
                        documents.append(doc)
            
            if not documents and not full_text_for_ocr_check.strip():
                 logger.warning(f"No text (direct or OCR) extracted from PDF '{original_filename}'.")

        logger.info(f"Created {len(documents)} chunks from PDF '{original_filename}'")
        return documents
    except Exception as e:
        logger.error(f"Error loading PDF {original_filename}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return []

def load_docx(file_path: str, file_id: str, original_filename: str) -> List[Document]:
    """Load DOCX file into documents with heading-aware chunking."""
    documents = []
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150, add_start_index=True)
    
    try:
        doc_obj = docx.Document(file_path)
        logger.info(f"DOCX '{original_filename}' loaded with {len(doc_obj.paragraphs)} paragraphs.")
        
        current_section_texts: List[str] = []
        # Default metadata for the first section or docs without explicit headings
        current_section_metadata: Dict[str, Any] = {
            "source": original_filename, "file_id": file_id, "file_type": "docx_section",
            "heading_level": 0, "heading_text": "General Content", "paragraph_start_index": 1,
            "doc_id_prefix": f"{file_id}_section_0"
        }
        section_count = 0

        for para_idx, para in enumerate(doc_obj.paragraphs):
            text = para.text.strip()
            # Basic heading detection (can be improved by checking specific style properties)
            is_heading = para.style.name.lower().startswith('heading')

            if is_heading and text: # New section starts
                if current_section_texts: # Process previous section
                    full_section_content = "\n".join(current_section_texts)
                    section_chunks = text_splitter.split_text(full_section_content)
                    for chunk_idx, chunk in enumerate(section_chunks):
                        chunk_meta = current_section_metadata.copy()
                        chunk_meta["chunk_in_section"] = chunk_idx + 1
                        chunk_meta["doc_id_prefix"] = f"{file_id}_section_{section_count}_chunk_{chunk_idx+1}"
                        documents.append(Document(page_content=chunk, metadata=chunk_meta))
                    current_section_texts = []
                
                section_count += 1
                current_section_metadata = {
                    "source": original_filename, "file_id": file_id, "file_type": "docx_section",
                    "heading_level": int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1, # Crude level
                    "heading_text": text, "paragraph_start_index": para_idx + 1,
                    "doc_id_prefix": f"{file_id}_section_{section_count}"
                }
                current_section_texts.append(text) # Include heading in its section
            elif text: # Regular paragraph
                current_section_texts.append(text)
        
        # Process the last accumulated section
        if current_section_texts:
            full_section_content = "\n".join(current_section_texts)
            section_chunks = text_splitter.split_text(full_section_content)
            for chunk_idx, chunk in enumerate(section_chunks):
                chunk_meta = current_section_metadata.copy()
                chunk_meta["chunk_in_section"] = chunk_idx + 1
                chunk_meta["doc_id_prefix"] = f"{file_id}_section_{section_count}_chunk_{chunk_idx+1}"
                documents.append(Document(page_content=chunk, metadata=chunk_meta))

        if not documents: # Fallback if no sections were created (e.g. empty doc or all empty paras)
            all_text_content = "\n".join([p.text.strip() for p in doc_obj.paragraphs if p.text.strip()])
            if all_text_content:
                doc_chunks = text_splitter.split_text(all_text_content)
                for i, chunk in enumerate(doc_chunks):
                    documents.append(Document(
                        page_content=chunk,
                        metadata={"source": original_filename, "file_id": file_id, 
                                  "file_type": "docx_full_chunk", "chunk_index": i+1,
                                  "doc_id_prefix": f"{file_id}_full_chunk_{i+1}"}
                    ))
        logger.info(f"Created {len(documents)} chunks from DOCX '{original_filename}'")
        return documents
    except Exception as e:
        logger.error(f"Error loading DOCX {original_filename}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return []

def load_pptx(file_path: str, file_id: str, original_filename: str) -> List[Document]:
    """Load PPTX file into documents, chunking content per slide."""
    documents = []
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=750, chunk_overlap=100, add_start_index=True)
    try:
        pres = pptx.Presentation(file_path)
        logger.info(f"PPTX '{original_filename}' has {len(pres.slides)} slides")
        for i, slide in enumerate(pres.slides):
            slide_texts_combined = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_texts_combined.append(shape.text.strip())
            
            if slide_texts_combined:
                slide_content = "\n\n".join(slide_texts_combined) 
                slide_chunks = text_splitter.split_text(slide_content)
                for chunk_idx, chunk in enumerate(slide_chunks):
                    documents.append(Document(
                        page_content=chunk,
                        metadata={
                            "slide": i + 1,
                            "chunk_on_slide": chunk_idx + 1,
                            "source": original_filename,
                            "file_id": file_id,
                            "file_type": "pptx_chunk",
                            "doc_id_prefix": f"{file_id}_slide_{i+1}_chunk_{chunk_idx+1}"
                        }
                    ))
        logger.info(f"Created {len(documents)} chunks from PPTX '{original_filename}'")
        return documents
    except Exception as e:
        logger.error(f"Error loading PPTX {original_filename}: {e}")
        return []

def load_text(file_path: str, file_id: str, original_filename: str) -> List[Document]:
    """Load TXT or MD file into documents using appropriate splitter."""
    documents = []
    file_ext = os.path.splitext(original_filename)[1].lower()
    
    if file_ext == '.md':
        splitter = MarkdownTextSplitter(chunk_size=1000, chunk_overlap=100)
        file_type_label = "md_chunk"
    else: # .txt or other plain text
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200, add_start_index=True)
        file_type_label = "text_chunk"
        
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        logger.info(f"Loaded text file '{original_filename}' with {len(content)} characters")
        
        if not content.strip():
            logger.warning(f"Text file {original_filename} is empty or contains only whitespace.")
            return []
        
        # Create chunks from the content
        text_chunks = splitter.split_text(content)
        for i, chunk in enumerate(text_chunks):
            doc = Document(
                page_content=chunk,
                metadata={
                    "chunk_index": i + 1,
                    "source": original_filename,
                    "file_id": file_id,
                    "file_type": file_type_label,
                    "doc_id_prefix": f"{file_id}_{file_type_label}_{i+1}"
                }
            )
            documents.append(doc)
        
        logger.info(f"Created {len(documents)} chunks from text file '{original_filename}'")
        return documents
    except Exception as e:
        logger.error(f"Error loading text file {original_filename}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return []

# --- File Processing Orchestration ---
def process_file(file_path: str, file_id: str, original_filename: str) -> Dict[str, Any]:
    """Process file based on its extension, load, chunk, and embed."""
    logger.info(f"Processing file: {original_filename} (id: {file_id}, path: {file_path})")
    
    if not os.path.exists(file_path):
        logger.error(f"File not found during processing: {file_path}")
        return {"success": False, "message": f"File not found: {original_filename}"}
    
    file_ext = os.path.splitext(original_filename)[1].lower()
    logger.info(f"File extension for {original_filename}: {file_ext}")
    
    docs: List[Document] = []
    file_type_processed: str = "unknown"

    if file_ext in ['.csv', '.xlsx', '.xls', '.xlsb']:
        docs = load_csv_xlsx(file_path, file_id, original_filename)
        file_type_processed = "tabular"
    elif file_ext == '.pdf':
        docs = load_pdf(file_path, file_id, original_filename)
        file_type_processed = "pdf"
    elif file_ext == '.docx':
        docs = load_docx(file_path, file_id, original_filename)
        file_type_processed = "docx"
    elif file_ext == '.pptx':
        docs = load_pptx(file_path, file_id, original_filename)
        file_type_processed = "pptx"
    elif file_ext in ['.txt', '.md']:
        docs = load_text(file_path, file_id, original_filename)
        file_type_processed = "text" if file_ext == '.txt' else "markdown"
    else:
        logger.error(f"Unsupported file type for {original_filename}: {file_ext}")
        # Clean up the uploaded file if unsupported
        try:
            os.remove(file_path)
            logger.info(f"Cleaned up unsupported file: {file_path}")
        except OSError as e_remove:
            logger.error(f"Error removing unsupported file {file_path}: {e_remove}")
        return {"success": False, "message": f"Unsupported file type: {file_ext}"}
    
    if not docs: # This means load_xxx returned an empty list
        logger.warning(f"No content extracted from {original_filename}. File might be empty or unreadable by current parsers.")
        # Clean up if no docs extracted, as it's not useful
        try:
            os.remove(file_path) # Remove the file from UPLOAD_DIR
            logger.info(f"Cleaned up file with no extracted content: {file_path}")
        except OSError as e_remove:
            logger.error(f"Error removing file with no content {file_path}: {e_remove}")
        return {
            "success": True, # File was "processed" (attempted), but yielded no docs
            "file_type": file_type_processed,
            "file_name": original_filename,
            "chunks_count": 0,
            "message": f"No content extracted or file empty: {original_filename}" # This message is key for the frontend error
        }
    
    logger.info(f"Embedding {len(docs)} documents from {original_filename}")
    try:
        chunks_count = embed(docs, file_id) 
        logger.info(f"Successfully embedded {chunks_count} chunks for {original_filename}")
        
        return {
            "success": True,
            "file_type": file_type_processed,
            "file_name": original_filename,
            "chunks_count": chunks_count,
            "file_id": file_id # Ensure file_id is returned
        }
    except Exception as e:
        logger.error(f"Error embedding documents for {original_filename}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return {"success": False, "message": f"Error embedding documents: {str(e)}"}

# --- API Endpoints ---

@app.post("/upload")
async def upload_file_endpoint(file: UploadFile = File(...)):
    """Upload and process a file. This endpoint handles the file saving and calls processing."""
    file_id = str(uuid.uuid4())
    original_filename = file.filename if file.filename else f"unknown_file_{file_id}"
    # Sanitize original_filename to prevent directory traversal or invalid characters
    original_filename = re.sub(r'[^\w\s.-]', '_', original_filename) 
    original_filename = original_filename[:200] # Limit length

    temp_file_name = f"{file_id}_{original_filename}" 
    file_path = os.path.join(UPLOAD_DIR, temp_file_name)
    
    try:
        logger.info(f"Received file upload: {original_filename}, assigning ID: {file_id}")
        
        with open(file_path, "wb") as f:
            content = await file.read()
            if not content:
                logger.warning(f"Uploaded file '{original_filename}' is empty.")
                try: os.remove(file_path) # Clean up the empty file created
                except OSError: pass
                raise HTTPException(status_code=400, detail="Uploaded file is empty.")
            f.write(content)
        
        logger.info(f"File {original_filename} saved to {file_path}, size: {len(content)} bytes. Processing...")
        
        result = process_file(file_path, file_id, original_filename)
        logger.info(f"File processing result for {original_filename} (ID: {file_id}): {result}")
        
        # Check if process_file indicated success but no content (which leads to the frontend error)
        if result.get("success") and result.get("chunks_count") == 0 and "No content extracted" in result.get("message", ""):
            # This is the scenario causing the frontend error.
            # The backend operation itself didn't fail catastrophically, but no data was usable.
            # The frontend will interpret this message.
            # We don't raise HTTPException here because the backend technically succeeded in its attempt.
            pass

        elif not result.get("success"):
            # If process_file itself failed (e.g., unsupported type, embedding error)
            raise HTTPException(status_code=500, detail=result.get("message", "File processing failed"))

        result["file_id"] = file_id 
        return result
            
    except HTTPException: 
        raise
    except Exception as e:
        logger.error(f"Error processing file upload for {original_filename}: {str(e)}", exc_info=True)
        if os.path.exists(file_path): 
            try: os.remove(file_path)
            except Exception as cleanup_e: logger.error(f"Error cleaning up {file_path} after error: {cleanup_e}")
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.post("/chat")
async def chat_endpoint(query: str = Form(...), history: Optional[str] = Form(None), file_id: Optional[str] = Form(None)):
    """Process a chat message with RAG."""
    try:
        logger.info(f"Received chat query: '{query}' for file_id: {file_id}")
        
        chat_history_parsed: List[Dict[str,str]] = []
        if history:
            import json
            try:
                chat_history_parsed = json.loads(history)
            except json.JSONDecodeError as e:
                logger.warning(f"Failed to parse chat history: {e}. Proceeding without history.")
        
        sql_indicators = ["list", "count", "how many", "average", "sum", "total", "compare", 
                          "sort", "filter", "show me", "table", "spreadsheet", "excel", 
                          "column", "row", "data for", "records of"]
        is_likely_sql_query = any(indicator in query.lower() for indicator in sql_indicators) or \
                              query.strip().lower().startswith("select ")
        
        retrieval_filter = {"file_id": file_id} if file_id else None
        retrieval_result = retrieve(query, metadata_filter=retrieval_filter)
        
        final_response_text = ""
        exportable_sql_query: Optional[str] = None
        context_for_llm = "" 

        if retrieval_result["type"] == "sql_result": 
            logger.info("Chat: Query was executed as direct SQL by RAG core's retrieve function.")
            context_for_llm = f"The SQL query `{retrieval_result['query']}` was executed.\n"
            if "columns" in retrieval_result and retrieval_result["columns"]:
                 context_for_llm += f"Columns: {', '.join(retrieval_result['columns'])}\n"
            if retrieval_result.get("data"):
                context_for_llm += f"Returned {len(retrieval_result['data'])} rows. Examples:\n"
                for i, row_item in enumerate(retrieval_result["data"][:3]): 
                    context_for_llm += f"- Row {i+1}: {row_item}\n"
            else:
                context_for_llm += "The query returned no data.\n"
            exportable_sql_query = retrieval_result['query']
        
        elif retrieval_result["type"] == "documents" and retrieval_result.get("data"):
            retrieved_docs = retrieval_result["data"]
            logger.info(f"Chat: Retrieved {len(retrieved_docs)} document chunks.")

            documents_are_tabular_related = any(
                doc.metadata.get("is_tabular") or doc.metadata.get("is_schema")
                for doc in retrieved_docs
            )

            if is_likely_sql_query and documents_are_tabular_related:
                logger.info("Chat: Query is likely SQL and docs are tabular. Attempting direct response generation.")
                direct_response_payload = process_query_for_direct_response(query, retrieved_docs, file_id)
                
                if direct_response_payload and direct_response_payload.get("type") == "sql_result" and direct_response_payload.get("direct_response_text"):
                    logger.info("Chat: Using direct response from response_generator.")
                    final_response_text = direct_response_payload["direct_response_text"]
                    exportable_sql_query = direct_response_payload.get("query")
                else:
                    logger.info("Chat: Direct response not applicable or failed. Building context for LLM from retrieved docs.")
            else: 
                logger.info("Chat: Not a direct SQL case or docs not tabular. Building context for LLM from retrieved docs.")
        
        elif retrieval_result["type"] == "error":
            logger.error(f"Chat: Error during document retrieval: {retrieval_result['message']}")
            final_response_text = f"I encountered an error trying to retrieve information: {retrieval_result['message']}. Please try again or check the uploaded document."
        else: 
            logger.info("Chat: No relevant documents found for the query.")
            final_response_text = "I couldn't find any relevant information in the current document(s) to answer your question."

        if not final_response_text: # If not handled by direct response or error
            if not context_for_llm: # Build context if not already set (e.g., by direct SQL exec summary)
                context_for_llm = "Retrieved context:\n"
                if retrieval_result["type"] == "documents" and retrieval_result.get("data"):
                    for doc_item in retrieval_result["data"]:
                        source_desc_parts = [f"Source: {doc_item.metadata.get('source', 'unknown')}"]
                        if doc_item.metadata.get('table_name'):
                            source_desc_parts.append(f"Table: {doc_item.metadata['table_name']}")
                            if doc_item.metadata.get('is_schema'): source_desc_parts.append("(Schema Info)")
                            elif doc_item.metadata.get('row_index'): source_desc_parts.append(f"Row approx: {doc_item.metadata['row_index']}")
                        elif doc_item.metadata.get('page'): source_desc_parts.append(f"Page approx: {doc_item.metadata['page']}")
                        elif doc_item.metadata.get('slide'): source_desc_parts.append(f"Slide approx: {doc_item.metadata['slide']}")
                        elif doc_item.metadata.get('heading_text'): source_desc_parts.append(f"Section: '{doc_item.metadata['heading_text']}'")
                        
                        internal_citation = f"InternalRef({', '.join(source_desc_parts)})" # For LLM's internal use
                        context_for_llm += f"Context from {internal_citation}:\n{doc_item.page_content}\n\n"
                else:
                    context_for_llm += "No specific context was retrieved for your query.\n"

            prompt_messages = [("system", SYSTEM_PROMPT)]
            for h_msg in chat_history_parsed: 
                role = "human" if h_msg.get("role") == "user" else "ai"
                prompt_messages.append((role, h_msg.get("content", "")))
            
            prompt_messages.append(("human", f"{context_for_llm}\n\nUser Question: {query}"))
            
            prompt = ChatPromptTemplate.from_messages(prompt_messages)
            chain = prompt | llm
            
            logger.info("Chat: Generating response from LLM...")
            llm_api_response = chain.invoke({}) 
            final_response_text = llm_api_response.content if hasattr(llm_api_response, 'content') else str(llm_api_response)
            logger.info("Chat: LLM response generated.")

        return {
            "response": final_response_text,
            "exportable": bool(exportable_sql_query),
            "sql": exportable_sql_query
        }
        
    except Exception as e:
        logger.error(f"Error processing chat query: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Error processing query: {str(e)}")

@app.get("/export")
async def export_data_endpoint(sql: str, filename: Optional[str] = None):
    """Export data as XLSX based on SQL query using the RAG core's SQL executor."""
    try:
        logger.info(f"Exporting data with SQL: {sql}")
        sql_result = execute_sql_query_rag(sql) 
        
        if sql_result.get("type") != "sql_result":
            error_message = sql_result.get('message', 'Unknown SQL error during export')
            logger.error(f"Invalid SQL query for export: {error_message}")
            raise HTTPException(status_code=400, detail=f"Invalid SQL query: {error_message}")
        
        data_to_export = sql_result.get("data", [])
        if not data_to_export:
            logger.warning("SQL query for export returned no data. Exporting empty Excel.")
            df = pd.DataFrame()
        else:
            df = pd.DataFrame(data_to_export)

        logger.info(f"Exporting {len(df)} rows.")
        
        export_filename_final = filename if filename else f"export_{uuid.uuid4().hex[:8]}.xlsx"
        if not export_filename_final.endswith(('.xlsx')):
            export_filename_final += '.xlsx'
        
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name="Sheet1")
        output.seek(0)
        
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename=\"{export_filename_final}\""}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error exporting data: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Error exporting data: {str(e)}")

@app.get("/tables")
async def list_tables_endpoint():
    """List available tables in the shared DuckDB database."""
    try:
        db_conn = get_shared_db_connection() 
        tables_df = db_conn.execute("SHOW TABLES").fetchdf()
        
        table_info_list = []
        if not tables_df.empty:
            table_names = list(tables_df['name'])
            logger.info(f"Available tables in DuckDB: {table_names}")
            for table_name_item in table_names:
                try:
                    count_result = db_conn.execute(f"SELECT COUNT(*) FROM \"{table_name_item}\"").fetchone()
                    count = count_result[0] if count_result else 0
                    
                    columns_df = db_conn.execute(f"PRAGMA table_info(\"{table_name_item}\")").fetchdf()
                    column_details = []
                    if not columns_df.empty:
                        column_details = [{"name": row['name'], "type": row['type']} for _, row in columns_df.iterrows()]
                    
                    table_info_list.append({
                        "name": table_name_item,
                        "rows": count,
                        "columns": column_details
                    })
                except Exception as e_table:
                    logger.error(f"Error getting info for table {table_name_item}: {e_table}")
                    table_info_list.append({
                        "name": table_name_item, "rows": "Error", "columns": [], "error": str(e_table)
                    })
        else:
            logger.info("No tables found in the shared DuckDB database.")
            
        return {"tables": table_info_list}
        
    except Exception as e:
        logger.error(f"Error listing tables: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Error listing tables: {str(e)}")

@app.get("/")
async def root_endpoint():
    return {"message": "DocuQuery API is running. Welcome!"}

if __name__ == "__main__":
    import uvicorn
    os.makedirs(UPLOAD_DIR, exist_ok=True) 
    logger.info(f"Starting Uvicorn server for DocuQuery API on port 8001.")
    uvicorn.run("main:app", host="0.0.0.0", port=8001, reload=True)