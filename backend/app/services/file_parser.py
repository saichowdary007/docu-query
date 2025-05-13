import csv
import os
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader
import markdown
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from typing import List, Dict, Any
from app.core.config import settings

# Ensure temp_uploads directory exists
os.makedirs(settings.TEMP_UPLOAD_FOLDER, exist_ok=True)

def parse_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)

def parse_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def parse_csv(file_path: str) -> pd.DataFrame:
    # For RAG, we might convert CSV rows to text or handle structured queries separately
    return pd.read_csv(file_path)

def parse_excel(file_path: str) -> Dict[str, pd.DataFrame]:
    # Returns a dictionary of sheet_name: dataframe
    return pd.read_excel(file_path, sheet_name=None)

def parse_md(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8') as f:
        return markdown.markdown(f.read())

def get_documents_from_file(file_path: str, filename: str) -> List[Document]:
    _, ext = os.path.splitext(filename.lower())
    content = ""
    metadata = {"source": filename}
    is_structured = False
    structured_data = None

    if ext == ".docx":
        content = parse_docx(file_path)
    elif ext == ".pptx":
        content = parse_pptx(file_path)
    elif ext == ".pdf":
        content = parse_pdf(file_path)
    elif ext == ".md":
        content = parse_md(file_path)
    elif ext == ".txt":
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    elif ext == ".csv":
        df = parse_csv(file_path)
        # For RAG: could convert rows to text, or store df for direct query
        # Option 1: Convert to text
        # content = df.to_string() 
        # Option 2: Store structured data reference
        is_structured = True
        structured_data = {"type": "csv", "data": df, "filename": filename}
        # Create a summary for embedding
        content = f"CSV file named {filename} with columns: {', '.join(df.columns.tolist())}. Data summary: {df.describe().to_string()}"

    elif ext in [".xls", ".xlsx"]:
        try:
            excel_data = parse_excel(file_path) # Dict of DFs
            is_structured = True
            structured_data = {"type": "excel", "data": excel_data, "filename": filename}
            
            # Create a summary for embedding
            content_parts = [f"Excel file named {filename} contains sheets: {', '.join(excel_data.keys())}."]
            
            for sheet_name, df in excel_data.items():
                # Convert column names to strings
                columns_str = ', '.join([str(col) for col in df.columns.tolist()])
                
                # Handle the data summary carefully to avoid type issues
                try:
                    # First try to convert the describe object to a string
                    summary_str = df.describe().to_string()
                except Exception as e:
                    # If that fails, try a safer approach
                    summary_str = f"[Summary data could not be stringified: {str(e)}]"
                
                content_parts.append(f"Sheet '{sheet_name}' has columns: {columns_str}. Data summary: {summary_str}")
            
            content = "\n".join(content_parts)
        except Exception as e:
            raise ValueError(f"Error processing Excel file: {str(e)}")
    else:
        # Potentially use Apache Tika for other formats or raise error
        raise ValueError(f"Unsupported file type: {ext}")

    if is_structured:
        # For structured data, we might create a single Document with summary,
        # and handle actual data operations separately.
        # The structured_data can be cached or stored for direct access.
        # For simplicity in RAG, we embed its summary.
        # For actual Pandas operations, data_handler.py will reload/use this.
        # A more robust system might store structured data in a proper DB.
        # For now, we'll rely on the content summary for retrieval, then use data_handler.
        # We can also store a reference to the original file path.
        metadata["structured_info"] = {"filename": filename, "type": structured_data["type"]} # type: ignore # type: "csv" or "excel"
        # Store the actual pandas DataFrame(s) in a simple cache or by filename reference
        # This part needs careful design for how data_handler accesses it.
        # For now, data_handler will re-load it if needed.

    if not content and not is_structured: # If content is empty and it's not a structured file processed differently
        return []
    
    # Chunk the content for RAG
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    texts = text_splitter.split_text(content)
    
    documents = [Document(page_content=text, metadata=metadata) for text in texts]
    return documents
