import csv
import os
from typing import Any, Dict, List

import markdown
import pandas as pd
from docx import Document as DocxDocument
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pypdf import PdfReader

from docuquery_ai.core.config import settings
from docuquery_ai.services.graph_service import get_knowledge_graph

# Ensure temp_uploads directory exists
os.makedirs(settings.TEMP_UPLOAD_FOLDER, exist_ok=True)


def parse_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs])


def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)


def parse_pdf(file_path: str) -> str:
    """
    Parse PDF and extract text with robust error handling and fallback methods.

    Args:
        file_path: Path to the PDF file

    Returns:
        Extracted text content from the PDF
    """
    try:
        # First try the standard method
        reader = PdfReader(file_path)

        # Check if the PDF is encrypted
        if reader.is_encrypted:
            try:
                # Try with empty password (some PDFs can be accessed this way)
                reader.decrypt("")
                print(f"Successfully decrypted PDF with empty password: {file_path}")
            except Exception as e:
                print(f"Cannot decrypt PDF {file_path}: {str(e)}")
                return f"[This PDF is encrypted and could not be processed: {os.path.basename(file_path)}]"

        # Extract text from each page with better error handling
        text = ""
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text() or ""
                text += page_text
                # If page is empty, add a note
                if not page_text.strip():
                    print(
                        f"Warning: Empty or non-text content on page {i+1} in {file_path}"
                    )
            except Exception as page_e:
                print(
                    f"Error extracting text from page {i+1} in {file_path}: {str(page_e)}"
                )
                text += f"\n[Error extracting text from page {i+1}]\n"

        # If we got no text at all, try a fallback method
        if not text.strip():
            print(
                f"Warning: No text extracted from {file_path}. Attempting fallback method."
            )
            # We could implement alternative extraction here if needed
            # e.g., using a different library or OCR for scanned PDFs
            text = f"[This document appears to contain no extractable text or may be a scanned PDF: {os.path.basename(file_path)}]"

        return text

    except Exception as e:
        print(f"Error parsing PDF {file_path}: {str(e)}")
        # Return a placeholder so the document isn't completely lost
        return f"[Error processing PDF document: {os.path.basename(file_path)}. Error: {str(e)}]"


def parse_csv(file_path: str) -> pd.DataFrame:
    # For RAG, we might convert CSV rows to text or handle structured queries separately
    return pd.read_csv(file_path)


def parse_excel(file_path: str) -> Dict[str, pd.DataFrame]:
    # Returns a dictionary of sheet_name: dataframe
    return pd.read_excel(file_path, sheet_name=None)


def parse_md(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return markdown.markdown(f.read())


def get_documents_from_file(file_path: str, filename: str) -> List[Document]:
    _, ext = os.path.splitext(filename.lower())
    content = ""
    metadata = {"source": filename}
    is_structured = False
    structured_data = None

    if ext == ".docx":
        content = parse_docx(file_path)
    elif ext == ".pptx":
        content = parse_pptx(file_path)
    elif ext == ".pdf":
        content = parse_pdf(file_path)
    elif ext == ".md":
        content = parse_md(file_path)
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
    elif ext == ".csv":
        df = parse_csv(file_path)
        # For RAG: could convert rows to text, or store df for direct query
        # Option 1: Convert to text
        # content = df.to_string()
        # Option 2: Store structured data reference
        is_structured = True
        structured_data = {"type": "csv", "data": df, "filename": filename}
        # Create a summary for embedding
        content = f"CSV file named {filename} with columns: {', '.join(df.columns.tolist())}. Data summary: {df.describe().to_string()}"

    elif ext in [".xls", ".xlsx"]:
        try:
            excel_data = parse_excel(file_path)  # Dict of DFs
            is_structured = True
            structured_data = {
                "type": "excel",
                "data": excel_data,
                "filename": filename,
            }

            # Create a summary for embedding
            content_parts = [
                f"Excel file named {filename} contains sheets: {', '.join(excel_data.keys())}."
            ]

            for sheet_name, df in excel_data.items():
                # Convert column names to strings
                columns_str = ", ".join([str(col) for col in df.columns.tolist()])

                # Handle the data summary carefully to avoid type issues
                try:
                    # First try to convert the describe object to a string
                    summary_str = df.describe().to_string()
                except Exception as e:
                    # If that fails, try a safer approach
                    summary_str = f"[Summary data could not be stringified: {str(e)}]"

                content_parts.append(
                    f"Sheet '{sheet_name}' has columns: {columns_str}. Data summary: {summary_str}"
                )

            content = "\n".join(content_parts)
        except Exception as e:
            raise ValueError(f"Error processing Excel file: {str(e)}")
    else:
        # Potentially use Apache Tika for other formats or raise error
        raise ValueError(f"Unsupported file type: {ext}")

    if is_structured:
        # For structured data, we might create a single Document with summary,
        # and handle actual data operations separately.
        # The structured_data can be cached or stored for direct access.
        # For simplicity in RAG, we embed its summary.
        # For actual Pandas operations, data_handler.py will reload/use this.
        # A more robust system might store structured data in a proper DB.
        # For now, we'll rely on the content summary for retrieval, then use data_handler.
        # We can also store a reference to the original file path.
        metadata["structured_info"] = {"filename": filename, "type": structured_data["type"]}  # type: ignore # type: "csv" or "excel"
        # Store the actual pandas DataFrame(s) in a simple cache or by filename reference
        # This part needs careful design for how data_handler accesses it.
        # For now, data_handler will re-load it if needed.

    if (
        not content and not is_structured
    ):  # If content is empty and it's not a structured file processed differently
        return []

    # Chunk the content for RAG
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    texts = text_splitter.split_text(content)

    documents = []
    knowledge_graph = get_knowledge_graph()
    for text in texts:
        doc_metadata = deepcopy(metadata)
        # Extract entities and relationships for potential graph integration
        extracted_graph_info = extract_entities_and_relationships(text, knowledge_graph)
        doc_metadata["graph_info"] = extracted_graph_info
        documents.append(Document(page_content=text, metadata=doc_metadata))
    save_knowledge_graph()
    return documents

def extract_entities_and_relationships(text: str, knowledge_graph) -> Dict[str, Any]:
    """
    Extracts entities and relationships from text and adds them to the knowledge graph.
    """
    extracted_info = {"entities": [], "relationships": []}

    # Simple keyword-based entity extraction for demonstration
    # In a real scenario, this would use an LLM or advanced NLP library
    entities = []
    if "project" in text.lower():
        entities.append(("project", "Concept"))
    if "report" in text.lower():
        entities.append(("report", "DocumentType"))
    if "user" in text.lower():
        entities.append(("user", "Person"))
    if "data" in text.lower():
        entities.append(("data", "Concept"))
    if "file" in text.lower():
        entities.append(("file", "Concept"))

    for entity_name, entity_type in entities:
        node_id = f"{entity_type.lower()}_{entity_name.lower()}"
        knowledge_graph.add_node(node_id, entity_type, {"name": entity_name})
        extracted_info["entities"].append({"id": node_id, "name": entity_name, "type": entity_type})

    # Simple keyword-based relationship extraction for demonstration
    # In a real scenario, this would use an LLM or advanced NLP library
    if "project" in text.lower() and "report" in text.lower():
        project_node_id = "concept_project"
        report_node_id = "documenttype_report"
        if knowledge_graph.get_node(project_node_id) and knowledge_graph.get_node(report_node_id):
            knowledge_graph.add_edge(report_node_id, project_node_id, "MENTIONS_PROJECT")
            extracted_info["relationships"].append({"source": report_node_id, "target": project_node_id, "type": "MENTIONS_PROJECT"})
    
    if "user" in text.lower() and "file" in text.lower():
        user_node_id = "person_user"
        file_node_id = "concept_file"
        if knowledge_graph.get_node(user_node_id) and knowledge_graph.get_node(file_node_id):
            knowledge_graph.add_edge(user_node_id, file_node_id, "UPLOADS")
            extracted_info["relationships"].append({"source": user_node_id, "target": file_node_id, "type": "UPLOADS"})

    return extracted_info
