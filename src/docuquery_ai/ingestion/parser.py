"""Utility functions for parsing different document formats.

The original implementation imported a number of optional third‑party
libraries at module import time.  This caused the package to fail to import
in minimal environments where those heavy dependencies were not installed.

To make the core package more lightweight and easier to test we now lazily
import those dependencies only when the corresponding parser function is
invoked.  Each parser raises a clear ``ImportError`` if the required library is
missing.
"""

from __future__ import annotations

import logging
import os
from typing import Dict

from docuquery_ai.core.config import settings

logger = logging.getLogger(__name__)

# Ensure temp_uploads directory exists
os.makedirs(settings.TEMP_UPLOAD_FOLDER, exist_ok=True)


def parse_docx(file_path: str) -> str:
    """Parse a DOCX file and extract its text content."""

    try:  # pragma: no cover - optional dependency
        from docx import Document as DocxDocument
    except Exception as exc:  # pragma: no cover - dependency not available
        raise ImportError("python-docx is required to parse DOCX files") from exc

    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs])


def parse_pptx(file_path: str) -> str:
    """Parse a PPTX file and extract text from all slides."""

    try:  # pragma: no cover - optional dependency
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - dependency not available
        raise ImportError("python-pptx is required to parse PPTX files") from exc

    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)


def parse_pdf(file_path: str) -> str:
    """Parse a PDF document and extract its textual content."""

    try:  # pragma: no cover - optional dependency
        from pypdf import PdfReader
    except Exception as exc:  # pragma: no cover - dependency not available
        raise ImportError("pypdf is required to parse PDF files") from exc

    try:
        reader = PdfReader(file_path)

        if reader.is_encrypted:
            try:
                reader.decrypt("")
                logger.warning("Successfully decrypted PDF with empty password: %s", file_path)
            except Exception as exc:
                logger.warning("Cannot decrypt PDF %s: %s", file_path, exc)
                return f"[This PDF is encrypted and could not be processed: {os.path.basename(file_path)}]"

        text = ""
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text() or ""
                text += page_text
                if not page_text.strip():
                    logger.warning("Empty or non-text content on page %s in %s", i + 1, file_path)
            except Exception as exc:
                logger.warning(
                    "Error extracting text from page %s in %s: %s", i + 1, file_path, exc
                )
                text += f"\n[Error extracting text from page {i+1}]\n"

        if not text.strip():
            logger.warning("No text extracted from %s. Attempting fallback method.", file_path)
            text = (
                f"[This document appears to contain no extractable text or may be a scanned PDF: {os.path.basename(file_path)}]"
            )

        return text
    except Exception as exc:
        logger.error("Error parsing PDF %s: %s", file_path, exc)
        return f"[Error processing PDF document: {os.path.basename(file_path)}. Error: {exc}]"


def parse_csv(file_path: str):
    """Parse a CSV file into a pandas :class:`DataFrame`."""

    import pandas as pd  # type: ignore  # pragma: no cover - optional dependency

    return pd.read_csv(file_path)


def parse_excel(file_path: str) -> Dict[str, "pd.DataFrame"]:
    """Parse an Excel workbook into a mapping of sheet names to DataFrames."""

    import pandas as pd  # type: ignore  # pragma: no cover - optional dependency

    return pd.read_excel(file_path, sheet_name=None)


def parse_md(file_path: str) -> str:
    """Parse a Markdown file and return its HTML representation."""

    try:  # pragma: no cover - optional dependency
        import markdown
    except Exception as exc:  # pragma: no cover - dependency not available
        raise ImportError("markdown package is required to parse Markdown files") from exc

    with open(file_path, "r", encoding="utf-8") as f:
        return markdown.markdown(f.read())

